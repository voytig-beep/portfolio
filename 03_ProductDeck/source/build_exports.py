import subprocess
from pathlib import Path

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
HTML = ROOT / "html" / "index.html"
EXPORTS = ROOT / "exports"
SLIDES = EXPORTS / "slides"
CONTACT = EXPORTS / "contact_sheet.jpg"
PPTX = ROOT / "SignalDesk_product_deck.pptx"
PDF = ROOT / "SignalDesk_product_deck.pdf"
GIF = EXPORTS / "SignalDesk_video_preview.gif"
PDF_SOURCE = EXPORTS / "pdf_source.html"


def edge_path() -> Path:
    candidates = [
        Path("C:/Program Files/Microsoft/Edge/Application/msedge.exe"),
        Path("C:/Program Files (x86)/Microsoft/Edge/Application/msedge.exe"),
        Path.home() / "AppData/Local/Microsoft/Edge/Application/msedge.exe",
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    raise SystemExit("Microsoft Edge executable not found")


def export_slides() -> list[Path]:
    edge = edge_path()
    SLIDES.mkdir(parents=True, exist_ok=True)
    result = []
    for index in range(1, 11):
        out = SLIDES / f"slide_{index:02}.png"
        url = f"{HTML.as_uri()}?export=1&slide={index}"
        subprocess.run(
            [
                str(edge),
                "--headless",
                "--disable-gpu",
                "--hide-scrollbars",
                "--force-device-scale-factor=1",
                "--run-all-compositor-stages-before-draw",
                "--window-size=1920,1080",
                f"--screenshot={out}",
                url,
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        result.append(out)
    return result


def build_pptx(images: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for image in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(PPTX)
    assert len(Presentation(PPTX).slides) == 10


def build_contact_sheet(images: list[Path]) -> None:
    thumb_w, thumb_h = 384, 216
    sheet = Image.new("RGB", (thumb_w * 2, thumb_h * 5), (8, 9, 10))
    for i, image_path in enumerate(images):
        img = Image.open(image_path).convert("RGB")
        img = ImageOps.fit(img, (thumb_w, thumb_h), method=Image.Resampling.LANCZOS)
        sheet.paste(img, ((i % 2) * thumb_w, (i // 2) * thumb_h))
    sheet.save(CONTACT, quality=92)


def build_pdf(images: list[Path]) -> None:
    pages = "\n".join(
        f'  <section class="page"><img src="slides/{image.name}" alt="Slide {i:02}"></section>'
        for i, image in enumerate(images, 1)
    )
    PDF_SOURCE.write_text(
        f"""<!doctype html>
<html lang="ru">
<head>
  <meta charset="utf-8">
  <style>
    @page {{ size: 16in 9in; margin: 0; }}
    html, body {{ margin: 0; padding: 0; background: #000; }}
    .page {{ width: 16in; height: 9in; page-break-after: always; break-after: page; overflow: hidden; }}
    .page:last-child {{ page-break-after: auto; break-after: auto; }}
    img {{ display: block; width: 100%; height: 100%; object-fit: cover; }}
  </style>
</head>
<body>
{pages}
</body>
</html>
""",
        encoding="utf-8",
    )
    subprocess.run(
        [
            str(edge_path()),
            "--headless",
            "--disable-gpu",
            "--run-all-compositor-stages-before-draw",
            f"--print-to-pdf={PDF}",
            "--print-to-pdf-no-header",
            PDF_SOURCE.as_uri(),
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


def build_gif(images: list[Path]) -> None:
    frames = []
    durations = []
    slides = [Image.open(path).convert("RGB").resize((960, 540), Image.Resampling.LANCZOS) for path in images]
    for i, current in enumerate(slides):
        frames.append(current)
        durations.append(850)
        if i == len(slides) - 1:
            continue
        nxt = slides[i + 1]
        for step in range(1, 4):
            frames.append(Image.blend(current, nxt, step / 4))
            durations.append(90)
    frames[0].save(GIF, save_all=True, append_images=frames[1:], duration=durations, loop=0, optimize=True)


def main() -> None:
    images = export_slides()
    if len(images) != 10:
        raise SystemExit(f"Expected 10 slides, got {len(images)}")
    build_pptx(images)
    build_contact_sheet(images)
    build_pdf(images)
    build_gif(images)
    for path in [PPTX, PDF, GIF, CONTACT]:
        print(path, path.stat().st_size)


if __name__ == "__main__":
    main()
