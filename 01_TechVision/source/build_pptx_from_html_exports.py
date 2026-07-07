from pathlib import Path

from PIL import Image, ImageOps, ImageDraw
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
SLIDES_DIR = ROOT / "exports" / "pptx_match_slides"
OUT = ROOT / "Noverra_коммерческое_предложение.pptx"
CONTACT = ROOT / "exports" / "pptx_match_contact_sheet.jpg"


def build_pptx() -> None:
    images = sorted(SLIDES_DIR.glob("slide_*.png"))
    if len(images) != 10:
        raise SystemExit(f"Expected 10 slide images, got {len(images)}")

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for image in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(OUT)

    check = Presentation(OUT)
    assert len(check.slides) == 10
    assert OUT.stat().st_size > 1_000_000


def build_contact_sheet() -> None:
    images = sorted(SLIDES_DIR.glob("slide_*.png"))
    thumb_w, thumb_h = 384, 216
    cols = 2
    rows = (len(images) + cols - 1) // cols
    sheet = Image.new("RGB", (thumb_w * cols, thumb_h * rows), (0, 0, 0))

    for i, image_path in enumerate(images):
        img = Image.open(image_path).convert("RGB")
        img = ImageOps.fit(img, (thumb_w, thumb_h), method=Image.Resampling.LANCZOS)
        draw = ImageDraw.Draw(img)
        draw.rectangle((0, 0, 46, 24), fill=(0, 0, 0))
        draw.text((8, 6), f"{i + 1:02}", fill=(255, 255, 255))
        sheet.paste(img, ((i % cols) * thumb_w, (i // cols) * thumb_h))

    CONTACT.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(CONTACT, quality=92)


if __name__ == "__main__":
    build_pptx()
    build_contact_sheet()
    print(OUT)
    print(CONTACT)
