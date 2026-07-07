from pathlib import Path

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
EXPORTS = ROOT / "exports" / "pptx_backgrounds"
OUT = ROOT / "TechVision_коммерческое_предложение.pptx"

BG = "030406"
PAPER = "F4F4F2"
TEXT = "F6F7F9"
MUTED = "B7BEC8"
INK = "121417"
INK_MUTED = "525963"
LINE = "FFFFFF"
CARD_DARK = "070A0F"
CARD_LIGHT = "FFFFFF"
ACCENT = "DCE7F5"

SLIDE_W = 13.333
SLIDE_H = 7.5


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def prepared_image(name: str, width: int, height: int) -> Path:
    EXPORTS.mkdir(parents=True, exist_ok=True)
    source = ASSETS / name
    target = EXPORTS / f"{Path(name).stem}_{width}x{height}.jpg"
    if target.exists() and target.stat().st_mtime >= source.stat().st_mtime:
        return target
    image = Image.open(source).convert("RGB")
    image = ImageOps.fit(image, (width, height), method=Image.Resampling.LANCZOS)
    image.save(target, quality=92)
    return target


def add_rect(slide, x, y, w, h, color, transparency=0, radius=False, line=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(LINE)
        shape.line.transparency = 78
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=24, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos Display"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_bg(slide, photo, mode="dark"):
    if mode == "dark":
        image = prepared_image(photo, 1920, 1080)
        slide.shapes.add_picture(str(image), Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
        add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, "000000", transparency=18)
        add_rect(slide, 0, 0, 7.6, SLIDE_H, "000000", transparency=9)
        return

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    if mode == "light-right":
        image = prepared_image(photo, 960, 1080)
        slide.shapes.add_picture(str(image), Inches(6.65), Inches(0), Inches(6.68), Inches(SLIDE_H))
        add_rect(slide, 5.85, 0, 1.3, SLIDE_H, PAPER, transparency=18)
    elif mode == "light-left":
        image = prepared_image(photo, 960, 1080)
        slide.shapes.add_picture(str(image), Inches(0), Inches(0), Inches(6.45), Inches(SLIDE_H))
        add_rect(slide, 5.95, 0, 1.2, SLIDE_H, PAPER, transparency=12)


def add_label(slide, label, mode):
    color = ACCENT if mode == "dark" else "6A7079"
    add_text(slide, label, .72, .48, 5.4, .26, 10, color, True)


def add_title(slide, title, lead, mode, x=.72, y=1.5, w=6.3):
    title_color = TEXT if mode == "dark" else INK
    lead_color = MUTED if mode == "dark" else INK_MUTED
    add_text(slide, title, x, y, w, 1.45, 36, title_color, True)
    add_text(slide, lead, x + .02, y + 1.68, w - .25, .82, 16, lead_color)


def add_card(slide, x, y, w, h, title, body, mode):
    if mode == "dark":
        add_rect(slide, x, y, w, h, CARD_DARK, transparency=28, radius=True, line=True)
        title_color, body_color = TEXT, MUTED
    else:
        add_rect(slide, x, y, w, h, CARD_LIGHT, transparency=8, radius=True, line=False)
        title_color, body_color = INK, INK_MUTED
    add_text(slide, title, x + .18, y + .18, w - .36, .34, 14, title_color, True)
    add_text(slide, body, x + .18, y + .64, w - .36, h - .72, 10.5, body_color)


def add_metric(slide, x, y, value, label, mode):
    add_card(slide, x, y, 1.55, 1.08, "", "", mode)
    color = TEXT if mode == "dark" else INK
    muted = MUTED if mode == "dark" else INK_MUTED
    add_text(slide, value, x + .18, y + .16, 1.15, .4, 28, color, True)
    add_text(slide, label, x + .18, y + .72, 1.16, .28, 8.5, muted)


def dark_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, data["photo"], "dark")
    add_label(slide, data["label"], "dark")
    add_title(slide, data["title"], data["lead"], "dark", w=6.2)

    cards = data.get("cards", [])
    if cards:
        if len(cards) == 4:
            positions = [(7.35, 1.42), (9.78, 1.42), (7.35, 3.0), (9.78, 3.0)]
            for (title, body), (x, y) in zip(cards, positions):
                add_card(slide, x, y, 2.18, 1.25, title, body, "dark")
        else:
            start_y = 1.72 if len(cards) == 3 else 1.45
            for i, (title, body) in enumerate(cards):
                add_card(slide, 7.7, start_y + i * 1.22, 4.45, 1.0, title, body, "dark")
    return slide


def light_slide(prs, data, mode="light-right"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, data["photo"], mode)
    add_label(slide, data["label"], mode)
    if mode == "light-left":
        text_x = 7.05
        card_x = 7.05
    else:
        text_x = .72
        card_x = .72
    add_title(slide, data["title"], data["lead"], "light", x=text_x, w=5.55)

    if "metrics" in data:
        for i, (value, label) in enumerate(data["metrics"]):
            add_metric(slide, .72 + i * 1.85, 5.05, value, label, "light")
    elif "cards" in data:
        cards = data["cards"]
        if len(cards) == 3:
            for i, (title, body) in enumerate(cards):
                add_card(slide, card_x, 4.35 + i * .88, 4.85, .72, title, body, "light")
        else:
            positions = [(card_x, 4.1), (card_x + 2.55, 4.1), (card_x, 5.32), (card_x + 2.55, 5.32)]
            for (title, body), (x, y) in zip(cards, positions):
                add_card(slide, x, y, 2.28, .94, title, body, "light")
    return slide


slides = [
    {
        "kind": "dark",
        "photo": "dark-workspace.jpg",
        "label": "01 / TECHVISION",
        "title": "Цифровые решения для бизнеса",
        "lead": "Коммерческое предложение для компаний, которым нужен понятный путь от идеи до работающего продукта.",
    },
    {
        "kind": "light-right",
        "photo": "office-glass.jpg",
        "label": "02 / О КОМПАНИИ",
        "title": "Запускаем цифровые продукты быстрее",
        "lead": "TechVision проектирует, разрабатывает и внедряет веб-сервисы, корпоративные платформы и аналитические инструменты.",
        "metrics": [("7", "лет на рынке"), ("80+", "проектов"), ("12", "отраслей")],
    },
    {
        "kind": "dark",
        "photo": "data-screen.jpg",
        "label": "03 / ПРОБЛЕМА",
        "title": "Технологии не должны тормозить рост",
        "lead": "Убираем узкие места, которые замедляют продажи, операции и клиентский опыт.",
        "cards": [
            ("Ручные процессы", "Команда тратит время на повторяющиеся операции."),
            ("Разрозненные данные", "Информация не помогает принимать решения."),
            ("Сложный клиентский путь", "Пользователи уходят там, где продукт должен продавать."),
        ],
    },
    {
        "kind": "light-left",
        "photo": "product-screen.jpg",
        "label": "04 / РЕШЕНИЕ",
        "title": "Собираем систему вокруг бизнес-целей",
        "lead": "Соединяем стратегию, UX, разработку и аналитику в один управляемый процесс.",
        "cards": [
            ("Структура", "Продуктовая логика до визуала и разработки."),
            ("Интерфейс", "Дизайн, который помогает двигаться дальше."),
            ("Запуск", "Рабочая система, готовая к развитию."),
        ],
    },
    {
        "kind": "light-right",
        "photo": "cover-laptop.jpg",
        "label": "05 / УСЛУГИ",
        "title": "Что мы можем сделать",
        "lead": "От первого MVP до корпоративной платформы с интеграциями и аналитикой.",
        "cards": [
            ("Корпоративные порталы", "Сайты, платформы и личные кабинеты."),
            ("CRM и внутренние системы", "Инструменты для продаж и операций."),
            ("Дашборды и аналитика", "Понятные данные для решений."),
            ("MVP и интеграции", "Быстрый запуск новых продуктов."),
        ],
    },
    {
        "kind": "dark",
        "photo": "premium-device.jpg",
        "label": "06 / ПРЕИМУЩЕСТВА",
        "title": "Не просто разработка. Управляемый результат.",
        "lead": "Делаем процесс прозрачным, а продукт готовым к росту после запуска.",
        "cards": [
            ("Прозрачные этапы", "Сроки, контрольные точки и ответственность."),
            ("Продуктовая экспертиза", "Смотрим на бизнес-результат."),
            ("Сильный дизайн", "Интерфейс помогает продавать."),
        ],
    },
    {
        "kind": "dark",
        "photo": "architecture.jpg",
        "label": "07 / КЕЙСЫ",
        "title": "Дизайн и технологии работают на бизнес",
        "lead": "Каждый проект связываем с измеримым эффектом для клиента.",
        "cards": [
            ("B2B-портал", "Минус 35% ручных операций."),
            ("Личный кабинет", "Плюс 22% к повторным оплатам."),
            ("Аналитический дашборд", "Отчёты за минуты вместо часов."),
        ],
    },
    {
        "kind": "dark",
        "photo": "dark-workspace.jpg",
        "label": "08 / ПРОЦЕСС",
        "title": "Понятный путь от первой встречи до запуска",
        "lead": "У проекта есть логика, поэтому клиент видит движение и результат.",
        "cards": [
            ("1 Диагностика", "Формулируем задачу и критерии успеха."),
            ("2 Концепция", "Определяем структуру и ключевые сценарии."),
            ("3 Дизайн", "Создаём UX/UI-прототип."),
            ("4 Запуск", "Разрабатываем, тестируем и поддерживаем."),
        ],
    },
    {
        "kind": "light-left",
        "photo": "team-work.jpg",
        "label": "09 / ФОРМАТЫ",
        "title": "Выберите формат под вашу задачу",
        "lead": "Можно начать с концепции, MVP или развития существующего продукта.",
        "cards": [
            ("Start", "Аудит, концепция и дорожная карта."),
            ("Build", "Дизайн и разработка MVP."),
            ("Scale", "Развитие, интеграции и аналитика."),
        ],
    },
    {
        "kind": "dark",
        "photo": "architecture.jpg",
        "label": "10 / КОНТАКТЫ",
        "title": "Готовы обсудить ваш цифровой продукт?",
        "lead": "Расскажите о задаче, и мы предложим оптимальный план запуска.",
        "cards": [
            ("hello@techvision.studio", "Почта для запроса коммерческого предложения."),
            ("techvision.studio", "Сайт компании и примеры проектов."),
        ],
    },
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    for data in slides:
        kind = data["kind"]
        if kind == "dark":
            dark_slide(prs, data)
        else:
            light_slide(prs, data, kind)

    prs.save(OUT)
    assert OUT.exists(), OUT
    assert len(prs.slides) == 10, len(prs.slides)
    print(OUT)


if __name__ == "__main__":
    build()
